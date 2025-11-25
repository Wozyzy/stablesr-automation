import os
import argparse
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

def parse_args():
    parser = argparse.ArgumentParser(description="Generate PPTX report from StableSR Seed Sweep.")
    parser.add_argument("--input_dir", type=str, required=True, help="Root directory containing fixed/random folders.")
    parser.add_argument("--output_pptx", type=str, default="seed_report.pptx", help="Path to save the output PPTX.")
    return parser.parse_args()

def parse_folder_name(folder_name):
    match = re.match(r"base(\d+)_x(\d+)p(\d+)", folder_name)
    if match:
        res = int(match.group(1))
        scale = float(f"{match.group(2)}.{match.group(3)}")
        return res, scale
    return None, None

def find_hq_image(folder_path):
    if not os.path.exists(folder_path):
        return None
    for f in os.listdir(folder_path):
        if f.lower().endswith(('.png', '.jpg', '.jpeg')) and "_lq" not in f:
            return os.path.join(folder_path, f)
    return None

def main():
    args = parse_args()
    
    data = {}
    modes = ["fixed", "random"]
    
    print(f"Scanning directory: {args.input_dir}")

    for mode in modes:
        mode_path = os.path.join(args.input_dir, mode)
        if not os.path.exists(mode_path):
            continue
            
        data[mode] = {}
        
        for folder_name in os.listdir(mode_path):
            folder_path = os.path.join(mode_path, folder_name)
            if not os.path.isdir(folder_path):
                continue
                
            res, scale = parse_folder_name(folder_name)
            if res is None:
                continue
            
            if res not in data[mode]:
                data[mode][res] = {}
            if scale not in data[mode][res]:
                data[mode][res][scale] = {}
                
            for run_folder in os.listdir(folder_path):
                if not run_folder.startswith("run_"):
                    continue
                try:
                    run_idx = int(run_folder.split("_")[1])
                except ValueError:
                    continue
                    
                run_path = os.path.join(folder_path, run_folder)
                hq_img = find_hq_image(run_path)
                if hq_img:
                    data[mode][res][scale][run_idx] = hq_img

    # Create Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "StableSR Seed Sweep Report"
    subtitle.text = "Fixed vs Random Seeds"

    blank_layout = prs.slide_layouts[6]
    
    margin_left = Inches(0.5)
    margin_top = Inches(1.5)
    gap = Inches(0.2)
    max_runs = 5
    
    available_width = prs.slide_width - (2 * margin_left)
    img_width = (available_width - ((max_runs - 1) * gap)) / max_runs

    for mode in modes:
        if mode not in data:
            continue
            
        sorted_res = sorted(data[mode].keys())
        for res in sorted_res:
            sorted_scales = sorted(data[mode][res].keys())
            for scale in sorted_scales:
                
                # New Slide for each (Mode, Res, Scale)
                slide = prs.slides.add_slide(blank_layout)
                
                # Title
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
                tf = txBox.text_frame
                tf.text = f"Mode: {mode.upper()} | Res: {res} | Scale: {scale}x"
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True
                
                for run_idx in range(max_runs):
                    img_path = data[mode][res][scale].get(run_idx)
                    
                    left = margin_left + (run_idx * (img_width + gap))
                    top = margin_top
                    
                    # Label (Run)
                    lblBox = slide.shapes.add_textbox(left, top - Inches(0.4), img_width, Inches(0.4))
                    lblTf = lblBox.text_frame
                    lblTf.text = f"Run {run_idx+1}"
                    lblTf.paragraphs[0].font.size = Pt(14)
                    
                    if img_path:
                        pic = slide.shapes.add_picture(img_path, left, top, width=img_width)
                        
                        # Dimensions
                        try:
                            with Image.open(img_path) as img:
                                w, h = img.size
                                dim_text = f"{w}x{h}"
                        except:
                            dim_text = "?x?"
                            
                        dimBox = slide.shapes.add_textbox(left, top + pic.height + Inches(0.1), img_width, Inches(0.4))
                        dimTf = dimBox.text_frame
                        dimTf.text = dim_text
                        dimTf.paragraphs[0].font.size = Pt(10)
                    else:
                        shape = slide.shapes.add_shape(
                            1, left, top, img_width, img_width
                        )
                        shape.text = "Missing"

    prs.save(args.output_pptx)
    print(f"PPTX Report saved to: {args.output_pptx}")

if __name__ == "__main__":
    main()
