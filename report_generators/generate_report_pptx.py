import os
import argparse
import re
from pptx import Presentation
from pptx.util import Inches, Pt

def parse_args():
    parser = argparse.ArgumentParser(description="Generate PPTX report from StableSR outputs (HQ Only).")
    parser.add_argument("--input_dir", type=str, required=True, help="Root directory containing the base{res}_x{scale} folders.")
    parser.add_argument("--output_pptx", type=str, default="report_hq.pptx", help="Path to save the output PPTX.")
    return parser.parse_args()

def parse_folder_name(folder_name):
    match = re.match(r"base(\d+)_x(\d+)p(\d+)", folder_name)
    if match:
        res = int(match.group(1))
        scale = float(f"{match.group(2)}.{match.group(3)}")
        return res, scale
    
    match = re.match(r"base(\d+)_x(\d+)", folder_name)
    if match:
        res = int(match.group(1))
        scale = float(match.group(2))
        return res, scale
        
    return None, None

def find_hq_image(folder_path):
    if not os.path.exists(folder_path):
        return None

    for f in os.listdir(folder_path):
        if not f.lower().endswith(('.png', '.jpg', '.jpeg')):
            continue
        if "_lq" in f:
            continue
        return os.path.join(folder_path, f)
    return None

def main():
    args = parse_args()
    
    data = {} # Key: scale, Value: {res: hq_path}
    resolutions = set()
    scales = set()

    print(f"Scanning directory: {args.input_dir}")
    
    if not os.path.exists(args.input_dir):
        print(f"Error: Directory '{args.input_dir}' does not exist.")
        return

    for folder_name in os.listdir(args.input_dir):
        folder_path = os.path.join(args.input_dir, folder_name)
        if not os.path.isdir(folder_path):
            continue
            
        res, scale = parse_folder_name(folder_name)
        if res is None:
            continue
            
        hq_path = find_hq_image(folder_path)
        
        if scale not in data:
            data[scale] = {}
        
        data[scale][res] = hq_path
        resolutions.add(res)
        scales.add(scale)

    sorted_resolutions = sorted(list(resolutions))
    sorted_scales = sorted(list(scales))
    
    print(f"Found Resolutions: {sorted_resolutions}")
    print(f"Found Scales: {sorted_scales}")
    
    if not sorted_scales:
        print("No matching folders found.")
        return

    # Create Presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio is default, usually 10x5.625 inches)
    # Let's use standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "StableSR Resolution Sweep Report"
    subtitle.text = "Comparison across Scales and Resolutions"

    # Content Slides: One per Scale
    blank_layout = prs.slide_layouts[6] # Blank
    
    # Layout calculations
    margin_left = Inches(0.5)
    margin_top = Inches(1.5)
    gap = Inches(0.2)
    
    # Available width for images
    available_width = prs.slide_width - (2 * margin_left)
    img_width = (available_width - ((len(sorted_resolutions) - 1) * gap)) / len(sorted_resolutions)
    
    for scale in sorted_scales:
        slide = prs.slides.add_slide(blank_layout)
        
        # Slide Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(1))
        tf = txBox.text_frame
        tf.text = f"Scale: {scale}x"
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        
        for i, res in enumerate(sorted_resolutions):
            hq_path = data[scale].get(res)
            
            left = margin_left + (i * (img_width + gap))
            top = margin_top
            
            # Label (Resolution)
            lblBox = slide.shapes.add_textbox(left, top - Inches(0.4), img_width, Inches(0.4))
            lblTf = lblBox.text_frame
            lblTf.text = f"Res: {res}"
            lblTf.paragraphs[0].font.size = Pt(18)
            lblTf.paragraphs[0].alignment = 2 # Center (PP_ALIGN.CENTER is not directly avail via enum in simple usage, but 2 usually works or import)
            # Actually, let's just leave alignment default or import
            
            if hq_path:
                # Add image
                # We add it with fixed width, height auto-scaled
                pic = slide.shapes.add_picture(hq_path, left, top, width=img_width)
                
                # Add dimensions label below
                # We need to read image size to display text
                # python-pptx doesn't give image size easily after add, but we can use PIL or assume
                # Let's use PIL just to get size for text
                try:
                    from PIL import Image
                    with Image.open(hq_path) as img:
                        w, h = img.size
                        dim_text = f"{w}x{h}"
                except:
                    dim_text = "?x?"
                    
                dimBox = slide.shapes.add_textbox(left, top + pic.height + Inches(0.1), img_width, Inches(0.4))
                dimTf = dimBox.text_frame
                dimTf.text = dim_text
                dimTf.paragraphs[0].font.size = Pt(12)
                
            else:
                # Placeholder for missing
                shape = slide.shapes.add_shape(
                    1, # msoShapeRectangle
                    left, top, img_width, img_width
                )
                shape.text = "Missing"

    prs.save(args.output_pptx)
    print(f"PPTX Report saved to: {args.output_pptx}")

if __name__ == "__main__":
    main()
